"""
Turns uploaded material (PDF, DOCX, PPTX, TXT) into clean, chunked text
ready for embedding. This is the input side of section 3 (Learning Material
Processing).
"""
from __future__ import annotations

import os
import re
from dataclasses import dataclass
from typing import List


@dataclass
class Chunk:
    text: str
    source: str          # filename
    location: str        # e.g. "page 4" or "slide 7" or "paragraph 12"
    chunk_id: str


def load_pdf(path: str) -> List[tuple[str, str]]:
    """Returns list of (text, location_label) per page."""
    try:
        import pymupdf as fitz  # PyMuPDF's newer import name
    except ImportError:
        import fitz  # older PyMuPDF versions
    doc = fitz.open(path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text")
        if text.strip():
            pages.append((text, f"page {i + 1}"))
    doc.close()
    return pages


def load_docx(path: str) -> List[tuple[str, str]]:
    import docx
    document = docx.Document(path)
    paras = []
    buf = []
    para_count = 0
    for p in document.paragraphs:
        if p.text.strip():
            buf.append(p.text)
        para_count += 1
        # group ~10 paragraphs per "location" to avoid a flood of tiny chunks
        if len(buf) >= 10:
            paras.append((" ".join(buf), f"paragraphs ~{para_count - 9}-{para_count}"))
            buf = []
    if buf:
        paras.append((" ".join(buf), f"paragraphs ~{para_count - len(buf) + 1}-{para_count}"))
    return paras


def load_pptx(path: str) -> List[tuple[str, str]]:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        texts.append(line)
        if texts:
            slides.append((" ".join(texts), f"slide {i + 1}"))
    return slides


def load_txt(path: str) -> List[tuple[str, str]]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    # split into ~500-word blocks for locations
    words = content.split()
    blocks = []
    step = 500
    for i in range(0, len(words), step):
        block = " ".join(words[i:i + step])
        blocks.append((block, f"section {i // step + 1}"))
    return blocks


LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".txt": load_txt,
    ".md": load_txt,
}


def clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"(\w)-\s+(\w)", r"\1\2", text)  # de-hyphenate line-wrapped words
    return text.strip()


def chunk_pages(pages: List[tuple[str, str]], source: str, max_chars: int = 1200, overlap: int = 150) -> List[Chunk]:
    """
    Chunk with overlap so a concept split across a page/slide boundary isn't
    lost to the retriever. Overlap matters more here than in generic RAG
    because textbook concepts often span boundaries.
    """
    chunks: List[Chunk] = []
    idx = 0
    for text, location in pages:
        text = clean_text(text)
        if not text:
            continue
        if len(text) <= max_chars:
            chunks.append(Chunk(text, source, location, f"{source}::{idx}"))
            idx += 1
            continue
        start = 0
        while start < len(text):
            end = min(start + max_chars, len(text))
            piece = text[start:end]
            chunks.append(Chunk(piece, source, location, f"{source}::{idx}"))
            idx += 1
            if end == len(text):
                break
            start = end - overlap
    return chunks


def ingest_file(path: str) -> List[Chunk]:
    ext = os.path.splitext(path)[1].lower()
    loader = LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {list(LOADERS)}")
    pages = loader(path)
    source = os.path.basename(path)
    return chunk_pages(pages, source)
